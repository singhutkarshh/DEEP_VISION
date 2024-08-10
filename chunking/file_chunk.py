import io
import copy
import re
from timeit import default_timer as timer
from io import BytesIO
from PIL import Image
from rag_tokenizer import tokenizer as rag_tokenizer
from parser.pdf_parser import RAGFlowPdfParser as PdfParser, PlainParser
from parser.ppt_parser import RAGFlowPptParser as PptParser
from parser.txt_parser import RAGFlowTxtParser as TxtParser
from parser.docx_parser import RAGFlowDocxParser as DocxParser
from PyPDF2 import PdfReader as pdf2_read
from docx import Document
from docx.image.exceptions import UnrecognizedImageError
from pptx import Presentation
from functools import reduce
from utils import concat_img, tokenize_table, naive_merge_docx, tokenize_chunks_docx
from parser.html_parser import RAGFlowHtmlParser as HtmlParser


def tokenize(d, t, eng):
    d["content_with_weight"] = t
    t = re.sub(r"</?(table|td|caption|tr|th)( [^<>]{0,12})?>", " ", t)
    d["content_ltks"] = rag_tokenizer.tokenize(t)
    d["content_sm_ltks"] = rag_tokenizer.fine_grained_tokenize(d["content_ltks"])
    
def is_english(texts):
    eng = 0
    if not texts: return False
    for t in texts:
        if re.match(r"[a-zA-Z]{2,}", t.strip()):
            eng += 1
    if eng / len(texts) > 0.8:
        return True
    return False

class Ppt(PptParser):  
    def __call__(self, fnm, from_page, to_page, callback=None):  
        txts = super().__call__(fnm, from_page, to_page)  
  
        callback(0.5, "Text extraction finished.")  
  
        imgs = []  
        print(fnm)  
        presentation = Presentation(fnm) 
        slides = list(presentation.slides)  
        for i, slide in enumerate(slides[from_page: to_page]):  
            for shape in slide.shapes:  
                if hasattr(shape, "image"):  
                    image = shape.image  
                    image_bytes = image.blob  
                    image_data = io.BytesIO(image_bytes)  
                    img = Image.open(image_data)  
                    imgs.append(img)  
  
        # assert len(imgs) == len(txts), "Slides text and image do not match: {} vs. {}".format(len(imgs), len(txts))  
        callback(0.9, "Image extraction finished")  
        self.is_english = is_english(txts)  
        return [(txts[i], imgs[i] if len(imgs) > i else None) for i in range(len(txts))] 

class Pdf(PdfParser):
    def __init__(self):
        super().__init__()

    def __garbage(self, txt):
        txt = txt.lower().strip()
        if re.match(r"[0-9\.,%/-]+$", txt):
            return True
        if len(txt) < 3:
            return True
        return False

    def __call__(self, filename, binary=None, from_page=0,
                 to_page=100000, zoomin=3, callback=None):
        callback(msg="OCR is running...")
        self.__images__(filename if not binary else binary,
                        zoomin, from_page, to_page, callback)
        callback(0.8, "Page {}~{}: OCR finished".format(
            from_page, min(to_page, self.total_page)))
        assert len(self.boxes) == len(self.page_images), "{} vs. {}".format(
            len(self.boxes), len(self.page_images))
        res = []
        for i in range(len(self.boxes)):
            lines = "\n".join([b["text"] for b in self.boxes[i]
                              if not self.__garbage(b["text"])])
            res.append((lines, self.page_images[i]))
        callback(0.9, "Page {}~{}: Parsing finished".format(
            from_page, min(to_page, self.total_page)))
        return res


class PlainPdf(PlainParser):
    def __call__(self, filename, binary=None, from_page=0,
                 to_page=100000, callback=None, **kwargs):
        self.pdf = pdf2_read(filename if not binary else BytesIO(binary))
        page_txt = []
        for page in self.pdf.pages[from_page: to_page]:
            page_txt.append(page.extract_text())
        callback(0.9, "Parsing finished")
        return [(txt, None) for txt in page_txt]

class Docx(DocxParser):
    def __init__(self):
        pass

    def get_picture(self, document, paragraph):
        img = paragraph._element.xpath('.//pic:pic')
        if not img:
            return None
        img = img[0]
        embed = img.xpath('.//a:blip/@r:embed')[0]
        related_part = document.part.related_parts[embed]
        try:
            image_blob = related_part.image.blob
        except UnrecognizedImageError:
            print(f"Unrecognized image format. Skipping image")
            return None
        try:
            image = Image.open(BytesIO(image_blob)).convert('RGB')
            return image
        except Exception as e:
            return None

    def __clean(self, line):
        line = re.sub(r"\u3000", " ", line).strip()
        return line

    def __call__(self, filename, binary=None, from_page=0, to_page=100000):
        self.doc = Document(
            filename) if not binary else Document(BytesIO(binary))
        pn = 0
        lines = []
        last_image = None
        for p in self.doc.paragraphs:
            if pn > to_page:
                break
            if from_page <= pn < to_page:
                current_image = None
                if p.text.strip():
                    if p.style.name == 'Caption':
                        former_image = None
                        if lines and lines[-1][1] and lines[-1][2] != 'Caption':
                            former_image = lines[-1][1].pop()
                        elif last_image:
                            former_image = last_image
                            last_image = None
                        lines.append((self.__clean(p.text), [former_image], p.style.name))
                    else:
                        current_image = self.get_picture(self.doc, p)
                        image_list = [current_image]
                        if last_image:
                            image_list.insert(0, last_image)
                            last_image = None
                        lines.append((self.__clean(p.text), image_list, p.style.name))
                else:
                    if current_image := self.get_picture(self.doc, p):
                        if lines:
                            lines[-1][1].append(current_image)
                        else:
                            last_image = current_image
            for run in p.runs:
                if 'lastRenderedPageBreak' in run._element.xml:
                    pn += 1
                    continue
                if 'w:br' in run._element.xml and 'type="page"' in run._element.xml:
                    pn += 1
        new_line = [(line[0], reduce(concat_img, line[1]) if line[1] else None) for line in lines]

        tbls = []
        for tb in self.doc.tables:
            html= "<table>"
            for r in tb.rows:
                html += "<tr>"
                i = 0
                while i < len(r.cells):
                    span = 1
                    c = r.cells[i]
                    for j in range(i+1, len(r.cells)):
                        if c.text == r.cells[j].text:
                            span += 1
                            i = j
                    i += 1
                    html += f"<td>{c.text}</td>" if span == 1 else f"<td colspan='{span}'>{c.text}</td>"
                html += "</tr>"
            html += "</table>"
            tbls.append(((None, html), ""))
        return new_line, tbls
    
def chunk(filename, binary=None, from_page=0, to_page=100000,
          lang="Chinese", callback=None, **kwargs):
    """
    The supported file formats are pdf, pptx.
    Every page will be treated as a chunk. And the thumbnail of every page will be stored.
    PPT file will be parsed by using this method automatically, setting-up for every PPT file is not necessary.
    """
    eng = lang.lower() == "english"
    
    parser_config = kwargs.get(
        "parser_config", {
            "chunk_token_num": 1280000, "delimiter": "", "layout_recognize": False})
     
    doc = {
        "docnm_kwd": filename,
        "title_tks": rag_tokenizer.tokenize(re.sub(r"\.[a-zA-Z]+$","", filename))
    }
    doc["title_sm_tks"] = rag_tokenizer.fine_grained_tokenize(doc["title_tks"])
    res = []
    if re.search(r"\.pptx?$", filename, re.IGNORECASE):
        pass
        ppt_parser = Ppt()
        for pn, (txt, img) in enumerate(ppt_parser(
                filename if not binary else binary, from_page, 1000000, callback)):
            d = copy.deepcopy(doc)
            pn += from_page
            d["image"] = img
            d["page_num_int"] = [pn + 1]
            d["top_int"] = [0]
            d["position_int"] = [(pn + 1, 0, img.size[0], 0, img.size[1])] if img else None
            tokenize(d, txt, eng)
            res.append(d)
        return res
    
    elif re.search(r"\.pdf$", filename, re.IGNORECASE):
        print("IN PDF")
        pdf_parser = Pdf() if kwargs.get(
            "parser_config", {}).get(
            "layout_recognize", True) else PlainPdf()
        for pn, (txt, img) in enumerate(pdf_parser(filename, binary,
                                                   from_page=from_page, to_page=to_page, callback=callback)):
            print("Processing page {}".format(pn + 1))
            callback(msg="Processing page {}".format(pn + 1))
            d = copy.deepcopy(doc)
            pn += from_page
            if img:
                d["image"] = img
            d["page_num_int"] = [pn + 1]
            d["top_int"] = [0]
            d["position_int"] = [
                (pn + 1, 0, img.size[0] if img else 0, 0, img.size[1] if img else 0)]
            tokenize(d, txt, eng)
            res.append(d)
        return res

    elif re.search(r"\.(txt|py|js|java|c|cpp|h|php|go|ts|sh|cs|kt)$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse.")
        sections = TxtParser()(filename,binary,parser_config.get("chunk_token_num", 128000))
        callback(0.8, "Finish parsing.")
        return sections
    elif re.search(r"\.docx$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse.")
        sections, tbls = Docx()(filename, binary)
        res = tokenize_table(tbls, doc, eng)    # just for table

        callback(0.8, "Finish parsing.")
        st = timer()

        chunks, images = naive_merge_docx(
            sections, int(parser_config.get(
                "chunk_token_num", 128)), parser_config.get(
                "delimiter", "\n!?。；！？"))

        if kwargs.get("section_only", False):
            return chunks

        res.extend(tokenize_chunks_docx(chunks, doc, eng, images))
        print("naive_merge({}): {}".format(filename, timer() - st))
        return res
    
    elif re.search(r"\.(htm|html)$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse.")
        sections = HtmlParser()(filename, binary)
        sections = [(l, "") for l in sections if l]
        callback(0.8, "Finish parsing.")
        return sections
    else:  
        raise NotImplementedError(
            "file type not supported yet(pptx, pdf supported)")